#!/usr/bin/env python3
"""
Build the lecture slide deck:
  "Practical LLMs & Agentic AI"
  Guest Lecture — University of Greenwich
  MSc Data Analytics & Management

Output: slides/practical-llms-agentic-ai-lecture.pptx
Import into Google Slides for final polish.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x7A, 0xCC)   # primary accent
ACCENT_TEAL  = RGBColor(0x00, 0xB4, 0xD8)   # secondary accent
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # success / metrics
ACCENT_ORANGE= RGBColor(0xE6, 0x7E, 0x22)   # warning / governance
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY     = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT     = RGBColor(0x33, 0x33, 0x33)
SECTION_BG   = RGBColor(0x0D, 0x47, 0xA1)   # section dividers

prs = Presentation()
prs.slide_width  = Inches(13.333)   # widescreen 16:9
prs.slide_height = Inches(7.5)

# ── Helper functions ────────────────────────────────────────────────

def add_bg(slide, color):
    """Fill slide background with solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    """Add a coloured rectangle behind everything."""
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def tx(shape, text, size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Set text on a shape's text frame (first paragraph)."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p

def title_slide(title, subtitle="", bg_color=DARK_BG):
    """Full-bleed title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_shape_bg(slide, bg_color)

    # Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2))
    tx(tb, title, size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.5))
        tx(tb2, subtitle, size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    return slide

def section_slide(part_num, title, bg_color=SECTION_BG):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    # Part label
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1))
    tx(tb, f"Part {part_num}", size=22, bold=False, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # Title
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.3), Inches(2))
    tx(tb2, title, size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    return slide

def content_slide(title, bullets, bg_color=DARK_BG, title_color=ACCENT_TEAL, bullet_color=WHITE, bullet_size=20):
    """Standard content slide with title and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    # Title bar accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), prs.slide_width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.8))
    tx(tb, title, size=30, bold=True, color=title_color, alignment=PP_ALIGN.LEFT)

    # Bullets
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = bullet_color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.level = 0
    return slide

def two_column_slide(title, left_title, left_bullets, right_title, right_bullets, bg_color=DARK_BG):
    """Two-column layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    # Title bar accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), prs.slide_width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.8))
    tx(tb, title, size=30, bold=True, color=ACCENT_TEAL)

    # Left column
    tb_lt = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.6))
    tx(tb_lt, left_title, size=22, bold=True, color=ACCENT_BLUE)

    tb_lb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = tb_lb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(6)

    # Right column
    tb_rt = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.6))
    tx(tb_rt, right_title, size=22, bold=True, color=ACCENT_BLUE)

    tb_rb = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
    tf2 = tb_rb.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(6)

    return slide

def table_slide(title, headers, rows, bg_color=DARK_BG):
    """Slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    # Title bar accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), prs.slide_width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.8))
    tx(tb, title, size=30, bold=True, color=ACCENT_TEAL)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(11.5)
    tbl_height = Inches(0.5) * n_rows
    left = Inches(0.9)
    top = Inches(1.4)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_width, tbl_height)
    table = table_shape.table

    col_width = int(tbl_width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_width

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x1A, 0x1A, 0x2E)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = LIGHT_GRAY
                p.font.name = "Calibri"

    return slide

def demo_slide(title, subtitle, url, key_points, business_takeaway, bg_color=DARK_BG):
    """Demo slide with screenshot placeholder and info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    # Title bar accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), prs.slide_width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.8))
    tf = tx(tb, f"DEMO: {title}", size=28, bold=True, color=ACCENT_ORANGE)

    # Subtitle
    tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.4))
    tx(tb_sub, subtitle, size=16, color=LIGHT_GRAY)

    # Screenshot placeholder (left half)
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(6.0), Inches(3.8)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    placeholder.line.color.rgb = ACCENT_BLUE
    placeholder.line.width = Pt(1)

    # Placeholder text
    ptf = placeholder.text_frame
    ptf.word_wrap = True
    p = ptf.paragraphs[0]
    p.text = "[Screenshot placeholder]"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_para(ptf, url, size=12, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    ptf.paragraphs[0].space_before = Pt(60)

    # Key points (right side)
    tb_pts = slide.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.3), Inches(3.0))
    tf2 = tb_pts.text_frame
    tf2.word_wrap = True
    for i, pt in enumerate(key_points):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"  {pt}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(6)

    # URL link
    tb_url = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(6.0), Inches(0.4))
    tx(tb_url, url, size=14, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # Business takeaway bar
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
    bar.line.fill.background()
    btf = bar.text_frame
    btf.word_wrap = True
    p = btf.paragraphs[0]
    p.text = f"Business Takeaway:  {business_takeaway}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    return slide

def quote_slide(quote, attribution, bg_color=DARK_BG):
    """Full-screen quote slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(3))
    tf = tx(tb, f'"{quote}"', size=32, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

    tb2 = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.8))
    tx(tb2, f"— {attribution}", size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    return slide

def metric_slide(title, metrics, bg_color=DARK_BG):
    """Big-number metrics slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, bg_color)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8))
    tx(tb, title, size=30, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(1.1), Inches(7.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    n = len(metrics)
    card_w = Inches(min(3.0, 11.0 / n))
    gap = (prs.slide_width - card_w * n) / (n + 1)

    for i, (number, label) in enumerate(metrics):
        left = int(gap * (i + 1) + card_w * i)
        top = Inches(1.8)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1)

        # Number
        tb_n = slide.shapes.add_textbox(left, Inches(2.5), card_w, Inches(1.5))
        tx(tb_n, number, size=48, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

        # Label
        tb_l = slide.shapes.add_textbox(left, Inches(4.0), card_w, Inches(1.5))
        tx(tb_l, label, size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# ════════════════════════════════════════════════════════════════════

# ── SLIDE 1: Title ──────────────────────────────────────────────────
s = title_slide(
    "Practical LLMs & Agentic AI",
    "Guest Lecture  |  University of Greenwich\nMSc Data Analytics & Management\n\nMondweep Chakravorty  &  Bence Csernak\nCo-Founders, Agentics Foundation — London Chapter"
)

# ── SLIDE 2: Agenda ─────────────────────────────────────────────────
content_slide("Today's Journey", [
    "Part 1    From ML to Agents  (20 min)",
    "Part 2    Architecture of Agent Systems  (20 min)",
    "Part 3    Live Demonstrations  (30-40 min)",
    "Part 4    Governance & Responsible AI  (20 min)",
    "Part 5    The Business Case & Future  (15 min)",
    "Part 6    Q&A & Discussion  (15 min)",
], bullet_size=22)

# ── SLIDE 3: About the speakers ─────────────────────────────────────
two_column_slide(
    "About Us",
    "Mondweep Chakravorty",
    [
        "Co-Founder, Agentics Foundation — London",
        "Background: Group Lotus, Jaguar Land Rover",
        "107 GitHub repos spanning agentic AI",
        "CISO London Summit presenter",
        "IC-ETSI 2025 — University of Greenwich",
    ],
    "Bence Csernak",
    [
        "Co-Founder, Agentics Foundation — London",
        "[Add Bence's bio details]",
    ],
)

# ── SLIDE 4: Your starting point ────────────────────────────────────
content_slide("You Already Know the Building Blocks", [
    "Machine Learning basics  — models learn from data",
    "Prediction  — models generate outputs from inputs",
    "Agency  — the theoretical concept of autonomous action",
    "",
    "Today we bridge from theory to the real world:",
    "What happens when you give a model tools, memory, and goals?",
])

# ════════════════════════════════════════════════════════════════════
# PART 1: FROM ML TO AGENTS
# ════════════════════════════════════════════════════════════════════
section_slide("1", "From ML to Agents")

# ── SLIDE 6: The spectrum ────────────────────────────────────────────
table_slide(
    "The AI Assistance Spectrum",
    ["Level", "Description", "Example", "Autonomy"],
    [
        ["Copilot", "Suggests completions", "Microsoft Copilot (your tool!)", "Low"],
        ["Assistant", "Follows instructions", "ChatGPT, Claude", "Medium"],
        ["Agent", "Pursues goals autonomously", "Claude Code, Forge agents", "High"],
        ["Swarm", "Coordinated multi-agent teams", "Claude-Flow, Agentic QE", "Very High"],
    ],
)

# ── SLIDE 7: What makes an agent ─────────────────────────────────────
content_slide("What Makes an AI System an Agent?", [
    "Tools  —  agents can read files, call APIs, run tests, edit code",
    "Memory  —  agents remember past outcomes and learn from them",
    "Goals  —  agents pursue objectives, not just respond to prompts",
    "Autonomy  —  agents decide their next action without step-by-step instruction",
    "Coordination  —  agents can communicate and collaborate with other agents",
    "",
    "Key distinction: reactive tools  vs  proactive agents that",
    '"accomplish goals, align stacked tasks, and complete them',
    'without direct supervision"',
])

# ── SLIDE 8: PACT Framework ─────────────────────────────────────────
two_column_slide(
    "The PACT Framework — Classifying Agentic AI",
    "The Four Dimensions",
    [
        "P — Proactive",
        "    Anticipate issues before they occur",
        "",
        "A — Autonomous",
        "    Self-operating with safety mechanisms",
        "",
        "C — Collaborative",
        "    Human expertise + AI capabilities",
        "",
        "T — Targeted",
        "    Risk-based focus on high-impact areas",
    ],
    "Think: SAE Levels for AI",
    [
        "Like autonomous vehicle levels (L0-L5),",
        "PACT classifies how autonomous an",
        "AI system is.",
        "",
        "Your Copilot experience?",
        "  Low P, Low A, Medium C, Low T",
        "",
        "A production agent swarm?",
        "  High P, High A, High C, High T",
    ],
)

# ── SLIDE 9: Copilot as entry point ─────────────────────────────────
content_slide("Your Copilot Is the Starting Point", [
    "Microsoft Copilot — you already use it under institutional licensing",
    "",
    "Copilot today:",
    "    Single-agent assistance  |  Reactive suggestions  |  Human-driven",
    "",
    "Where the industry is going:",
    "    Multi-agent orchestration  |  Proactive quality  |  Graduated autonomy",
    "",
    "The same principles (tools, memory, goals) apply at every scale",
    "Today we'll show you what's possible when you scale up",
])

# ════════════════════════════════════════════════════════════════════
# PART 2: ARCHITECTURE OF AGENT SYSTEMS
# ════════════════════════════════════════════════════════════════════
section_slide("2", "Architecture of Agent Systems")

# ── SLIDE 11: Agent anatomy ──────────────────────────────────────────
content_slide("How Agents Work — No Code Required", [
    "Perception  —  Agents read code, APIs, test results, user requirements",
    "Reasoning   —  LLMs process context and decide next actions",
    "Action      —  Agents execute through tools (file editing, API calls, tests)",
    "Memory      —  Agents learn from outcomes and store patterns for reuse",
    "Coordination — Agents communicate via shared memory or direct protocols",
    "",
    'Think of it as a team of analysts:',
    'each member has a speciality, shares a project dashboard,',
    'and coordinates through stand-ups and handoffs.',
])

# ── SLIDE 12: Claude-Flow ────────────────────────────────────────────
content_slide("Real Architecture: Claude-Flow", [
    "The leading open-source multi-agent orchestration platform",
    "",
    "5-Layer Architecture:",
    "  1. Entry Layer — CLI and MCP servers with security hardening",
    "  2. Routing Layer — Q-Learning routers, 8 mixture-of-experts, 42+ skills",
    "  3. Swarm Coordination — hierarchical & peer-to-peer topologies",
    "  4. Agent Pool — 60+ specialized agents (coders, testers, reviewers...)",
    "  5. Intelligence Layer — self-optimization, vector search, 9 RL algorithms",
], bullet_size=18)

# ── SLIDE 13: Claude-Flow metrics ────────────────────────────────────
metric_slide("Claude-Flow in Numbers", [
    ("500K", "Downloads\nworldwide"),
    ("60+", "Specialized\nAgents"),
    ("84.8%", "SWE-Bench\nSolve Rate"),
    ("2.8-4.4x", "Faster Task\nCompletion"),
])

# ── SLIDE 14: Forge specialization ───────────────────────────────────
content_slide("Agent Specialization — Forge", [
    "Forge: autonomous quality engineering with 8 specialized agents",
    "",
    "  Specification Verifier  — checks requirements (Sonnet)",
    "  Test Runner  — executes test suites (Haiku)",
    "  Failure Analyzer  — diagnoses what went wrong (Sonnet)",
    "  Bug Fixer  — writes the fix (Opus)",
    "  Quality Gate Enforcer  — checks all gates pass (Haiku)",
    "  Accessibility Auditor  — WCAG compliance (Sonnet)",
    "  Auto-Committer  — commits passing code (Haiku)",
    "  Learning Optimizer  — improves future runs (Sonnet)",
    "",
    "Like a project team — each member matched to the right skill level",
], bullet_size=17)

# ── SLIDE 15: Forge loop ─────────────────────────────────────────────
content_slide("The Autonomous Quality Loop", [
    "Specify  ->  Test  ->  Analyze  ->  Fix  ->  Audit  ->  Gate  ->  Commit  ->  Learn",
    "",
    "This loop runs continuously without human intervention:",
    "  Quality is forged in, not bolted on",
    "",
    "But humans remain in the loop for:",
    "  Strategy and design decisions",
    "  Reviewing critical changes",
    "  Adjusting confidence thresholds",
    "  Overriding agent decisions when needed",
])

# ════════════════════════════════════════════════════════════════════
# PART 3: LIVE DEMONSTRATIONS
# ════════════════════════════════════════════════════════════════════
section_slide("3", "Live Demonstrations")

# ── SLIDE 17: Demo overview ──────────────────────────────────────────
table_slide(
    "Demo Lineup",
    ["#", "Demo", "Key Metric", "Theme"],
    [
        ["D1", "40-Minute App Build (Pre-Route)", "5 agents, 40 min to working app", "Speed to market"],
        ["D2", "AI Security Audit (MAESTRO)", "23 vulns in 4 hrs, 16:1 ROI", "Governance at speed"],
        ["D3", "AI Analytics (Auto-Analyst)", "Real-time pricing intelligence", "Your domain transformed"],
        ["D4", "Knowledge Graph (LBS)", "3,963 nodes enriched for $14", "AI at minimal cost"],
        ["D5", "WASM Image Filters (Live)", "10-50x WASM vs JS speedup", "Privacy by design"],
        ["D6", "Assam AI Governance (Live)", "20-30 days to 5-7 days", "Government transparency"],
        ["D7", "Driftwise (Live)", "Voice-first, GPS-triggered", "Voice-first AI"],
    ],
)

# ── SLIDE 18-19: Demo D1 — Pre-Route ────────────────────────────────
demo_slide(
    "40-Minute App Build — Pre-Route",
    "London Agentics Meetup  |  August 2025",
    "github.com/mondweep/london-agentics-meetup-13aug-25",
    [
        "5 specialized agents collaborated:",
        "  Coordinator, Research, Backend,",
        "  Frontend, Test Engineer",
        "",
        "Built a complete traffic monitoring",
        "app for Kent residents in 40 minutes",
        "",
        "Express.js + TypeScript backend",
        "with mock Google Maps/TomTom APIs",
    ],
    "What would take weeks, done in under an hour — with quality built in",
)

# ── SLIDE 20-21: Demo D2 — MAESTRO Security ─────────────────────────
demo_slide(
    "AI Security Audit — MAESTRO Framework",
    "CISO London Summit  |  Rela8 Group",
    "github.com/mondweep/agentic-ai-security-demo-rela8group-ciso-london-summit",
    [
        "Analysed ElizaOS: 127 files, 42K+ lines",
        "Found 23 vulnerabilities in 4 hours",
        "  (vs 15-20 days traditional audit)",
        "",
        "3 Critical findings:",
        "  Plugin Supply Chain Attack (56/70)",
        "  scored highest risk",
        "",
        "Cost: $5-15K vs $50-150K traditional",
    ],
    "Governance and security can be accelerated, not just development",
)

metric_slide("Security Audit ROI", [
    ("96%", "Faster\n4 hrs vs 15-20 days"),
    ("90%", "Cheaper\n$5-15K vs $50-150K"),
    ("16:1", "3-Year ROI\nmitigation value"),
    ("23", "Vulnerabilities\nfound automatically"),
])

# ── SLIDE 22: Demo D3 — Auto-Analyst ────────────────────────────────
demo_slide(
    "AI-Powered Analytics — Auto-Analyst",
    "AI transforming traditional data analytics workflows",
    "github.com/mondweep/Auto-Analyst",
    [
        "Automotive dealership analytics",
        "platform with AI-powered insights",
        "",
        "Interactive dashboards",
        "Pricing intelligence",
        "Market opportunity analysis",
        "",
        "This is YOUR domain —",
        "data analytics transformed by agents",
    ],
    "AI transforms traditional analytical workflows — and data analysts are best positioned to lead",
)

# ── SLIDE 23: Demo D4 — Knowledge Graph ─────────────────────────────
demo_slide(
    "University Knowledge Graph — LBS",
    "AI-driven content intelligence at minimal cost",
    "github.com/mondweep/university-pitch",
    [
        "Semantic enrichment of london.edu",
        "3,963 nodes, 3,953 edges",
        "",
        "Phases: Crawling -> Domain Modeling",
        "  -> Semantic Enrichment",
        "",
        "Sentiment analysis, topic extraction,",
        "persona classification",
        "",
        "Total cost: ~$14",
    ],
    "AI applied to a university — like Greenwich. Near-zero marginal cost for content intelligence",
)

# ── SLIDE 24: Demo D5 — WASM Image Filters ──────────────────────────
demo_slide(
    "WASM Image Filters",
    "Privacy-first processing — your data never leaves the browser",
    "https://wasm-tinkering.netlify.app/",
    [
        "Upload an image, apply 10+ filters",
        "JS vs WASM benchmark side-by-side",
        "",
        "10-50x performance improvement",
        "with Rust/WASM over JavaScript",
        "",
        "ALL processing runs locally —",
        "data never leaves the device",
        "",
        "Built with agentic AI tools",
    ],
    "Privacy-by-design and performance optimization are not mutually exclusive",
)

# ── SLIDE 25: Demo D6 — LuitPlayer ──────────────────────────────────
demo_slide(
    "LuitPlayer — Grand Piano",
    "7-agent swarm built a 7-octave piano with sheet music OMR",
    "https://grand-piano-thisismon.netlify.app/",
    [
        "88 keys (C1-B7), 5 timbres",
        "Built by 7-agent Claude Flow swarm:",
        "  Coordinator, Frontend, Audio,",
        "  OMR, Test, Deploy, Integration",
        "",
        "Optical Music Recognition for",
        "Assamese sheet music (Tesseract WASM)",
        "",
        "Multi-touch, sustain pedal, mixer",
    ],
    "Multi-agent collaboration can produce sophisticated creative applications",
)

# ── SLIDE 26: Demo D7 — Kumno ───────────────────────────────────────
demo_slide(
    "Kumno — Khasi Travel Companion",
    "AI translation enabling cultural preservation and tourism",
    "https://kumno.netlify.app/",
    [
        "80+ Khasi phrases with phonetic guides",
        "8 categories: Greetings, Taxis,",
        "  Markets, Food, Directions...",
        "",
        "Bah/Kong honorific toggle",
        "  (cultural sensitivity built in)",
        "",
        "AI-powered real-time translation",
        "Fully offline-capable PWA",
    ],
    "AI translation bridging language barriers while preserving cultural nuance",
)

# ── SLIDE 27: Demo D8 — Assam AI Governance ─────────────────────────
demo_slide(
    "Assam AI Governance",
    "AI for government transparency and anti-fraud",
    "https://assam-ai-usecase-governance-pwa.netlify.app/",
    [
        "Property Registration Digitalization:",
        "  20-30 days reduced to 5-7 days",
        "  80% remote completion target",
        "",
        "Infrastructure Cost Auditing:",
        "  AI anomaly detection (Isolation Forest)",
        "  Catches 3x cost inflation",
        "  50+ crore INR annual savings target",
        "",
        "Trilingual: English, Assamese, Hindi",
    ],
    "AI applied directly to government governance and transparency — measurable cost savings",
)

# ── SLIDE 28: Demo D9 — Driftwise ───────────────────────────────────
demo_slide(
    "Driftwise — Discover Interesting Facts",
    "Voice-first, driver-safe AI that transforms commutes into discovery",
    "https://driftwise-discover-interesting-facts.netlify.app/",
    [
        "GPS-triggered historical fact delivery",
        "No user interaction required",
        "",
        "Voice-first design — no screen glancing",
        "5-second voice command window",
        "",
        "Anti-generic filtering: rejects clichés,",
        "surfaces named individuals, specific",
        "dates, unusual events",
        "",
        "Gemini Live API for real-time voice",
    ],
    "Voice-first design pattern — responsible AI with driver safety as the primary constraint",
)

# ════════════════════════════════════════════════════════════════════
# PART 4: GOVERNANCE & RESPONSIBLE AI
# ════════════════════════════════════════════════════════════════════
section_slide("4", "Governance & Responsible AI")

# ── SLIDE 30: Honest limitations ─────────────────────────────────────
content_slide("Let's Be Honest — Current Limitations", [
    "Weak business-rule understanding",
    "    Agents can miss domain-specific nuances that humans catch intuitively",
    "",
    "Model reliability issues",
    "    LLMs can hallucinate, lose context, or produce inconsistent outputs",
    "",
    "Over-engineering tendencies",
    "    Agents sometimes add unnecessary complexity",
    "",
    "Severity calibration challenges",
    "    Distinguishing critical issues from minor ones remains difficult",
    "",
    "Real maintenance overhead",
    "    Agent systems require monitoring, tuning, and governance",
], bullet_size=17)

quote_slide(
    "These are engineering trade-offs, not showstoppers.\nCritical thinking — not hype — is what governance requires.",
    "The honest assessment approach"
)

# ── SLIDE 32: MAESTRO framework ──────────────────────────────────────
table_slide(
    "MAESTRO — 7-Layer Security Framework for Agentic AI",
    ["Layer", "Focus", "Example Concerns"],
    [
        ["1. Foundational Models", "LLM integrations, inference", "Prompt security, model poisoning"],
        ["2. Data Operations", "Databases, RAG, vectors", "Data leakage, poisoned embeddings"],
        ["3. Agent Frameworks", "Orchestration, decision logic", "Rogue agents, decision manipulation"],
        ["4. Deployment & Infra", "Runtime, APIs, networking", "Unauthorized access, API abuse"],
        ["5. Evaluations", "Logging, monitoring, testing", "Blind spots, inadequate audit trails"],
        ["6. Security & Compliance", "Auth, secrets, policies", "Credential theft, policy violations"],
        ["7. Agent Ecosystem", "Plugins, actions, tools", "Supply chain attacks, malicious plugins"],
    ],
)

# ── SLIDE 33: Quality Gates ──────────────────────────────────────────
content_slide("Quality Gates — Guardrails for Autonomous Agents", [
    "7 blocking gates from Forge that agents must pass:",
    "",
    "  1. Functional Correctness — code compiles and passes unit tests",
    "  2. Behavioural Compliance — Gherkin specifications all pass",
    "  3. Code Coverage — meets minimum thresholds (85% baseline, 95% critical)",
    "  4. Security Violations — no SAST/DAST findings above threshold",
    "  5. API Contract Validation — responses match expected contracts",
    "  6. Accessibility — WCAG AA compliance",
    "  7. Resilience — chaos testing passes",
    "",
    "Agents operate freely within these boundaries",
    "Gates are the governance mechanism, not humans watching every step",
], bullet_size=17)

# ── SLIDE 34: Confidence tiers ───────────────────────────────────────
content_slide("Confidence Tiers — Trust Earned Through Evidence", [
    "How do you know when to trust an agent with more autonomy?",
    "",
    "Platinum (10+ successes)  —  fully autonomous operation",
    "Gold (3-9 successes)  —  autonomous with monitoring",
    "Silver (1-2 successes)  —  requires human review",
    "",
    "Trust is earned incrementally through empirical evidence",
    "  Not assigned by role or assumption",
    "",
    "Like vendor performance ratings in procurement —",
    "  you wouldn't give a new vendor a £1M contract on day one",
])

# ── SLIDE 35: TinyDancer cost governance ─────────────────────────────
table_slide(
    "TinyDancer — Cost Governance Through Model Routing",
    ["Complexity", "Model", "Use Cases", "Cost"],
    [
        ["0-20 (Simple)", "Haiku", "Syntax fixes, type additions", "Lowest"],
        ["20-70 (Standard)", "Sonnet", "Bug fixes, test generation", "Moderate"],
        ["70-100 (Complex)", "Opus", "Architecture, security analysis", "Highest"],
    ],
)

metric_slide("Cost Governance Impact", [
    ("75%", "Token Cost\nReduction"),
    ("3", "Model Tiers\nIntelligent routing"),
    ("85%", "Multi-Provider\nCost Savings"),
])

# ── SLIDE 37: Human in the loop ──────────────────────────────────────
content_slide("Human-in-the-Loop — A Design Principle, Not a Limitation", [
    "Agents augment, never fully replace human judgement",
    "",
    "Humans own:",
    "    Strategy & design decisions",
    "    Critical business-rule validation",
    "    Ethical oversight & escalation",
    "    Adjusting confidence thresholds",
    "",
    "Agents own:",
    "    Execution at speed and scale",
    "    Pattern recognition across large codebases",
    "    Continuous quality monitoring",
    "    Routine decision-making within guardrails",
])

# ── SLIDE 38: Discussion prompt ──────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, RGBColor(0x0D, 0x47, 0xA1))
tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.8))
tx(tb, "Discussion", size=30, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

tb2 = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(5.0))
tf = tx(tb2, '"If you were deploying an agentic AI system in your organisation,\nwhat quality gates would you define?"', size=26, color=WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "", size=14, color=WHITE)
add_para(tf, '"How would you measure whether an agent has earned\nenough trust for increased autonomy?"', size=26, color=WHITE, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# PART 5: BUSINESS CASE & FUTURE
# ════════════════════════════════════════════════════════════════════
section_slide("5", "The Business Case & Future")

# ── SLIDE 40: ROI table ──────────────────────────────────────────────
table_slide(
    "ROI Evidence — Numbers You Can Take to a Board Meeting",
    ["Application", "Traditional", "AI-Driven", "Improvement"],
    [
        ["Security Audit", "15-20 days, $50-150K", "4 hours, $5-15K", "96% faster, 90% cheaper"],
        ["Task Completion", "Baseline", "2.8-4.4x faster", "Up to 340% improvement"],
        ["Token Costs", "Baseline", "75% reduction", "Via model routing"],
        ["Content Enrichment", "Manual", "~$14 for 3,963 nodes", "Near-zero marginal cost"],
        ["App Development", "Weeks/months", "40 minutes (Pre-Route)", "Radical speed to market"],
        ["Test Coverage", "Manual QA", "85-95% automated", "Quality gates enforce standards"],
    ],
)

# ── SLIDE 41: Cross-industry ─────────────────────────────────────────
two_column_slide(
    "Cross-Industry Applications",
    "Already Happening",
    [
        "Healthcare — diagnostic AI, care coordination",
        "Automotive — pricing intelligence, inventory",
        "Education — recruitment, Socratic tutors",
        "Public Sector — fraud prevention, governance",
        "Security — automated vulnerability assessment",
        "Media — content discovery, sentiment analysis",
    ],
    "Where It's Going",
    [
        "Single agents -> coordinated swarms",
        "Reactive testing -> proactive quality",
        "Manual governance -> automated compliance",
        "Human-only decisions -> graduated autonomy",
        "Agent-Ready Web (ARW) — websites",
        "  designed for agent consumption",
        "A2A protocols — agents from different",
        "  vendors collaborating",
    ],
)

# ── SLIDE 42: Higher Ed transformation ───────────────────────────────
content_slide("Agentic AI Is Already Transforming Universities", [
    "Recruitment  —  24/7 digital concierges handling credit evaluations and campus tours",
    "Academic Support  —  Socratic AI tutors generating practice problems in real time",
    "Predictive Intervention  —  identifying struggling students before midterms",
    "Mental Health  —  triage agents offering coping strategies, escalating serious cases",
    "Finance  —  automated invoice processing, grant identification, proposal drafting",
    "Compliance  —  regulatory reporting generating audit-ready documents",
    "",
    '"Institutions operationalising AI early will widen their performance gap"',
    "— Inside Higher Ed, January 2026",
], bullet_size=17)

# ── SLIDE 43: Your role ─────────────────────────────────────────────
content_slide("Your Role as Data Analysts", [
    "You are uniquely positioned for this revolution",
    "",
    "Your superpowers:",
    "    Metrics & measurement — you know how to evaluate effectiveness",
    "    Critical analysis — you can spot when AI claims don't hold up",
    "    Business context — you bridge technical capability and business value",
    "    Governance thinking — data governance skills transfer directly",
    "",
    "The governance challenge needs business-minded people,",
    "  not just engineers",
    "",
    "Data analysts will evaluate, govern, and propose agentic AI systems",
])

# ════════════════════════════════════════════════════════════════════
# PART 6: Q&A & DISCUSSION
# ════════════════════════════════════════════════════════════════════
section_slide("6", "Q&A & Discussion")

# ── SLIDE 45: Discussion topics ──────────────────────────────────────
content_slide("Discussion Topics", [
    "1.  What business processes in your industry could benefit from agentic AI?",
    "",
    "2.  How would you design governance for an agent that manages financial data?",
    "",
    "3.  What ethical considerations arise when agents learn from outcomes?",
    "",
    "4.  How does your experience with Microsoft Copilot map to the PACT framework?",
    "",
    "5.  If you could deploy one agent in your workplace tomorrow, what would it do?",
], bullet_size=20)

# ── SLIDE 46: Glossary ──────────────────────────────────────────────
table_slide(
    "Glossary — Quick Reference",
    ["Term", "Business-Friendly Definition"],
    [
        ["Agent", "AI system that pursues goals autonomously — a digital team member"],
        ["Swarm", "Coordinated team of AI agents — a project team with specialized roles"],
        ["Quality Gate", "Checkpoint that work must pass — like approval stages in a process"],
        ["Confidence Tier", "Trust level based on track record — like vendor performance ratings"],
        ["Model Routing", "Directing tasks to right-sized AI — assigning work by complexity"],
        ["PACT", "Framework for classifying agent autonomy: Proactive, Autonomous, Collaborative, Targeted"],
        ["MAESTRO", "7-layer security framework for evaluating agentic AI systems"],
        ["MCP", "Model Context Protocol — standard way AI agents connect with tools (USB for AI)"],
    ],
)

# ── SLIDE 47: Resources ─────────────────────────────────────────────
two_column_slide(
    "Resources & Further Reading",
    "Open-Source Repos",
    [
        "Claude-Flow — github.com/ruvnet/claude-flow",
        "Forge — github.com/ikennaokpala/forge",
        "Agentic QE — github.com/proffesor-for-testing/agentic-qe",
        "Pre-Route — github.com/mondweep/london-agentics-meetup-13aug-25",
        "MAESTRO Demo — github.com/mondweep/agentic-ai-security-demo...",
        "Auto-Analyst — github.com/mondweep/Auto-Analyst",
        "University KG — github.com/mondweep/university-pitch",
    ],
    "Live Demo Apps",
    [
        "WASM Filters — wasm-tinkering.netlify.app",
        "LuitPlayer — grand-piano-thisismon.netlify.app",
        "Kumno — kumno.netlify.app",
        "Assam Gov — assam-ai-usecase-governance-pwa.netlify.app",
        "Driftwise — driftwise-discover-interesting-facts.netlify.app",
        "",
        "Frameworks",
        "PACT — forge-quality.dev",
        "Agentic QE — agentic-qe.dev",
    ],
)

# ── SLIDE 48: Agentics Foundation ────────────────────────────────────
content_slide("Agentics Foundation — London Chapter", [
    "Co-founded by Mondweep Chakravorty and Bence Csernak",
    "",
    "Global chapters: London, Serbia, India (Assam), New Zealand",
    "",
    "Partnership with University of Greenwich — IC-ETSI 2025 conference",
    "",
    "Regular meetups with hands-on workshops on Agentic Engineering Practices",
    "",
    "Hackathons proving the viability of agent-assisted rapid development",
    "",
    "Get involved: join our London meetups and hackathons",
])

# ── SLIDE 49: Thank you ─────────────────────────────────────────────
s = title_slide(
    "Thank You",
    "Mondweep Chakravorty  &  Bence Csernak\nCo-Founders, Agentics Foundation — London Chapter\n\ngithub.com/mondweep\n\nQuestions?"
)

# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════

output_path = os.path.join(os.path.dirname(__file__), "practical-llms-agentic-ai-lecture.pptx")
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to {output_path}")
