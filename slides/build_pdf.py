#!/usr/bin/env python3
"""
Generate a PDF version of the lecture slides using ReportLab.
Reads the PPTX and renders each slide as a landscape PDF page.
"""

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch, mm
from reportlab.lib.colors import Color
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os

PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "practical-llms-agentic-ai-lecture.pptx")
PDF_PATH  = os.path.join(os.path.dirname(os.path.abspath(__file__)), "practical-llms-agentic-ai-lecture.pdf")

# Register fonts
for name, path in [
    ("DejaVu", "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ("DejaVu-Bold", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
]:
    if os.path.exists(path):
        pdfmetrics.registerFont(TTFont(name, path))

FONT = "DejaVu"
FONT_B = "DejaVu-Bold"

prs = Presentation(PPTX_PATH)

# Page dimensions (points) — match slide aspect ratio 13.333 x 7.5 inches
PAGE_W = 13.333 * inch
PAGE_H = 7.5 * inch

# Conversion helpers
def emu2pt(emu):
    """EMU to points (1 inch = 914400 EMU = 72 pt)."""
    return (emu / 914400.0) * 72.0

def rgb_from_pptx(rgb_obj):
    """Convert pptx RGBColor to reportlab Color."""
    try:
        s = str(rgb_obj)
        r = int(s[0:2], 16) / 255.0
        g = int(s[2:4], 16) / 255.0
        b = int(s[4:6], 16) / 255.0
        return Color(r, g, b)
    except:
        return Color(1, 1, 1)

def get_fill(shape):
    try:
        f = shape.fill
        if f.type is not None and f.fore_color and f.fore_color.rgb:
            return rgb_from_pptx(f.fore_color.rgb)
    except:
        pass
    return None

def get_font_color(run):
    try:
        if run.font.color and run.font.color.rgb:
            return rgb_from_pptx(run.font.color.rgb)
    except:
        pass
    return Color(1, 1, 1)

def get_font_size(run, default=12):
    try:
        if run.font.size:
            return run.font.size / 12700.0  # EMU to pt
    except:
        pass
    return default

def is_bold(run):
    try:
        return run.font.bold is True
    except:
        return False


c = canvas.Canvas(PDF_PATH, pagesize=(PAGE_W, PAGE_H))

for slide_idx, slide in enumerate(prs.slides):
    # Background: default dark
    c.setFillColor(Color(26/255, 26/255, 46/255))
    c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)

    # Sort shapes by z-order (draw order)
    shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

    for shape in shapes:
        left = emu2pt(shape.left) if shape.left else 0
        top_from_top = emu2pt(shape.top) if shape.top else 0
        w = emu2pt(shape.width) if shape.width else 0
        h = emu2pt(shape.height) if shape.height else 0

        # ReportLab y is from bottom
        bot = PAGE_H - top_from_top - h

        # Draw filled shapes
        if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, 1, 5):
            fill_c = get_fill(shape)
            if fill_c:
                c.setFillColor(fill_c)
                c.rect(left, bot, w, h, fill=1, stroke=0)

        # Draw tables
        if shape.has_table:
            table = shape.table
            n_cols = len(table.columns)
            n_rows = len(table.rows)
            col_ws = [emu2pt(table.columns[j].width) for j in range(n_cols)]
            row_h = h / n_rows if n_rows else 12

            for i, row in enumerate(table.rows):
                x = left
                cell_top = top_from_top + i * row_h
                cell_bot = PAGE_H - cell_top - row_h

                for j, cell in enumerate(row.cells):
                    cw = col_ws[j]

                    # Cell fill
                    try:
                        if cell.fill and cell.fill.type is not None:
                            fc = cell.fill.fore_color
                            if fc and fc.rgb:
                                c.setFillColor(rgb_from_pptx(fc.rgb))
                                c.rect(x, cell_bot, cw, row_h, fill=1, stroke=0)
                    except:
                        pass

                    # Cell border
                    c.setStrokeColor(Color(0.25, 0.25, 0.35))
                    c.setLineWidth(0.3)
                    c.rect(x, cell_bot, cw, row_h, fill=0, stroke=1)

                    # Cell text
                    txt = cell.text.strip()
                    if txt:
                        fsz = 7.5 if i == 0 else 6.5
                        fc = Color(1, 1, 1) if i == 0 else Color(0.88, 0.88, 0.88)
                        fn = FONT_B if i == 0 else FONT

                        c.setFont(fn, fsz)
                        c.setFillColor(fc)

                        # Simple text truncation to fit cell
                        max_chars = int(cw / (fsz * 0.45))
                        lines = []
                        words = txt.split()
                        line = ""
                        for word in words:
                            if len(line) + len(word) + 1 <= max_chars:
                                line = line + " " + word if line else word
                            else:
                                lines.append(line)
                                line = word
                        if line:
                            lines.append(line)

                        ty = cell_bot + row_h - fsz - 1.5
                        for ln in lines[:4]:  # max 4 lines per cell
                            c.drawString(x + 3, ty, ln)
                            ty -= fsz + 1

                    x += cw

        # Draw text frames
        if shape.has_text_frame:
            tf = shape.text_frame
            y_cursor = top_from_top + 2  # start near top of textbox

            for para in tf.paragraphs:
                full_text = para.text
                if not full_text.strip():
                    y_cursor += 6
                    continue

                # Get formatting
                if para.runs:
                    run = para.runs[0]
                    color = get_font_color(run)
                    fsz = get_font_size(run, 12)
                    bold = is_bold(run)
                else:
                    color = Color(1, 1, 1)
                    fsz = 12
                    bold = False

                fsz = max(6, min(fsz, 36))
                fn = FONT_B if bold else FONT

                c.setFont(fn, fsz)
                c.setFillColor(color)

                # Alignment
                align = 'L'
                try:
                    if para.alignment == PP_ALIGN.CENTER:
                        align = 'C'
                    elif para.alignment == PP_ALIGN.RIGHT:
                        align = 'R'
                except:
                    pass

                # Word wrap
                usable_w = w - 6 if w > 10 else PAGE_W - left - 10
                max_chars = max(10, int(usable_w / (fsz * 0.48)))

                lines = []
                words = full_text.split()
                line = ""
                for word in words:
                    if len(line) + len(word) + 1 <= max_chars:
                        line = line + " " + word if line else word
                    else:
                        lines.append(line)
                        line = word
                if line:
                    lines.append(line)

                for ln in lines:
                    y_pt = PAGE_H - y_cursor  # flip to bottom-origin
                    if y_pt < 10:
                        break

                    if align == 'C':
                        tw = c.stringWidth(ln, fn, fsz)
                        x_draw = left + (w - tw) / 2
                    elif align == 'R':
                        tw = c.stringWidth(ln, fn, fsz)
                        x_draw = left + w - tw - 3
                    else:
                        x_draw = left + 3

                    c.drawString(max(x_draw, 2), y_pt, ln)
                    y_cursor += fsz + 2

                y_cursor += 2

    # Slide number
    c.setFont(FONT, 7)
    c.setFillColor(Color(0.47, 0.47, 0.55))
    c.drawRightString(PAGE_W - 15, 10, str(slide_idx + 1))

    c.showPage()

c.save()
print(f"Saved {len(prs.slides)} slides to {PDF_PATH}")
print(f"File size: {os.path.getsize(PDF_PATH) / 1024:.0f} KB")
